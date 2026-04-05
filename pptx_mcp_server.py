"""
pptx_mcp_server.py
------------------
MCP Server #1 — Presentation Tools

Uses FastMCP with @mcp.tool() decorator (clean, modern style).

Tools:
  - create_presentation   : Initialise a new .pptx file
  - add_slide             : Add a styled slide (title / content / summary)
  - write_text            : Overwrite text on an existing slide
  - add_image_placeholder : Add a labelled image placeholder box
  - save_presentation     : Save the .pptx file to disk
  - get_slide_count       : Return how many slides exist so far

Run:
  python pptx_mcp_server.py
"""

# ── MCP ───────────────────────────────────────────────────────────────────────
from mcp.server.fastmcp import FastMCP

# ── Standard library ─────────────────────────────────────────────────────────
import os

# ── python-pptx ───────────────────────────────────────────────────────────────
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── typing ────────────────────────────────────────────────────────────────────
from typing import Optional, List

# ──────────────────────────────────────────────────────────────────────────────
# FastMCP server instance
# ──────────────────────────────────────────────────────────────────────────────
mcp = FastMCP("pptx-mcp-server")

# ──────────────────────────────────────────────────────────────────────────────
# Global presentation state (one .pptx per session)
# ──────────────────────────────────────────────────────────────────────────────
_prs: Optional[Presentation] = None
_output_path: str = "output.pptx"

# ──────────────────────────────────────────────────────────────────────────────
# Design constants  —  "Midnight Star" theme
# ──────────────────────────────────────────────────────────────────────────────
PRIMARY   = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
SECONDARY = RGBColor(0x1B, 0x26, 0x3B)   # lighter navy
ACCENT    = RGBColor(0xFF, 0xD7, 0x00)   # golden yellow
TEXT_MAIN = RGBColor(0xFF, 0xFF, 0xFF)   # white
TEXT_MUTE = RGBColor(0xCA, 0xDC, 0xFC)   # ice-blue

# ──────────────────────────────────────────────────────────────────────────────
# Internal drawing helpers  (not exposed as MCP tools)
# ──────────────────────────────────────────────────────────────────────────────

def _apply_bg(slide, color: RGBColor) -> None:
    """Fill the slide background with a solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, text: str, left, top, width, height,
                 font_size: int = 14, bold: bool = False,
                 color: RGBColor = TEXT_MAIN,
                 align=PP_ALIGN.LEFT, italic: bool = False) -> None:
    """Add a text box with consistent styling."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def _draw_title_slide(slide, title: str, subtitle: str = "", slide_index: int = 0) -> None:
    """Render a title / opening slide with premium styling."""
    _apply_bg(slide, PRIMARY)

    # 1. Slide Index Badge (e.g., SLIDE 1)
    _add_textbox(slide, f"SLIDE {slide_index + 1}",
                 Inches(0.4), Inches(0.3), Inches(2.0), Inches(0.3),
                 font_size=10, color=TEXT_MUTE, bold=True)

    # 2. Type Badge (e.g., TITLE)
    badge_bg = slide.shapes.add_shape(5, Inches(0.4), Inches(0.7), Inches(1.2), Inches(0.35))
    badge_bg.fill.solid()
    badge_bg.fill.fore_color.rgb = ACCENT
    badge_bg.line.fill.background()
    
    _add_textbox(slide, "   TITLE", 
                 Inches(0.4), Inches(0.72), Inches(1.2), Inches(0.3),
                 font_size=9, bold=True, color=PRIMARY)

    # Left accent bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(2.5), Inches(0.15), Inches(2.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    # Title text
    _add_textbox(slide, title,
                 Inches(0.4), Inches(2.6), Inches(9.2), Inches(1.6),
                 font_size=40, bold=True, color=TEXT_MAIN)

    # Subtitle text
    if subtitle:
        _add_textbox(slide, subtitle,
                     Inches(0.4), Inches(4.3), Inches(9.2), Inches(0.8),
                     font_size=18, color=TEXT_MUTE)


def _draw_content_slide(slide, title: str, bullets: List, slide_index: int = 0, placeholder_label: str = "[Image: Relevant Visual]") -> None:
    """Render a premium split-layout content slide with text (L) and image placeholder (R)."""
    _apply_bg(slide, PRIMARY)

    # 1. Slide Index Badge (e.g., SLIDE 2)
    _add_textbox(slide, f"SLIDE {slide_index + 1}",
                 Inches(0.4), Inches(0.3), Inches(2.0), Inches(0.3),
                 font_size=10, color=TEXT_MUTE, bold=True)

    # 2. Type Badge (e.g., CONTENT)
    badge_bg = slide.shapes.add_shape(5, Inches(0.4), Inches(0.7), Inches(1.2), Inches(0.35))
    badge_bg.fill.solid()
    badge_bg.fill.fore_color.rgb = SECONDARY
    badge_bg.line.fill.background()
    
    _add_textbox(slide, "  CONTENT", 
                 Inches(0.4), Inches(0.72), Inches(1.2), Inches(0.3),
                 font_size=9, bold=True, color=TEXT_MUTE)

    # 3. Main Title (Reduced width for split layout)
    _add_textbox(slide, title,
                 Inches(0.4), Inches(1.3), Inches(5.0), Inches(1.0),
                 font_size=28, bold=True, color=TEXT_MAIN)

    # 4. Bullet points matching the '▸' style ── (restricted width)
    top = Inches(2.4)
    for bullet in bullets[:5]:
        _add_textbox(slide, "▸", 
                     Inches(0.4), top, Inches(0.3), Inches(0.5),
                     font_size=14, color=ACCENT)
        
        _add_textbox(slide, bullet,
                     Inches(0.7), top + Inches(0.015), Inches(4.5), Inches(0.8),
                     font_size=15, color=TEXT_MUTE)
        top += Inches(0.8)

    # 5. Right-side Image Placeholder Card
    ph_bg = slide.shapes.add_shape(1, Inches(5.6), Inches(1.2), Inches(4.0), Inches(5.0))
    ph_bg.fill.solid()
    ph_bg.fill.fore_color.rgb = RGBColor(0x1B, 0x26, 0x3B) # Match secondary
    ph_bg.line.color.rgb = ACCENT
    ph_bg.line.width = Pt(1.5)

    _add_textbox(slide, str(placeholder_label),
                 Inches(5.7), Inches(3.4), Inches(3.8), Inches(0.6),
                 font_size=12, color=TEXT_MUTE, italic=True, align=PP_ALIGN.CENTER)


def _draw_summary_slide(slide, title: str, points: List, slide_index: int = 0) -> None:
    """Render a closing summary slide matching the frontend preview exactly."""
    _apply_bg(slide, PRIMARY)

    # 1. Slide Index Badge (e.g., SLIDE 5)
    _add_textbox(slide, f"SLIDE {slide_index + 1}",
                 Inches(0.4), Inches(0.3), Inches(2.0), Inches(0.3),
                 font_size=10, color=TEXT_MUTE, bold=True)

    # 2. Type Badge (e.g., SUMMARY)
    badge_bg = slide.shapes.add_shape(5, Inches(0.4), Inches(0.7), Inches(1.2), Inches(0.35))
    badge_bg.fill.solid()
    badge_bg.fill.fore_color.rgb = RGBColor(0x98, 0xFF, 0x98) # Mint green
    badge_bg.line.fill.background()
    
    _add_textbox(slide, " SUMMARY", 
                 Inches(0.4), Inches(0.72), Inches(1.2), Inches(0.3),
                 font_size=9, bold=True, color=PRIMARY)

    # 3. Main Title (Left-aligned to match preview)
    _add_textbox(slide, title,
                 Inches(0.4), Inches(1.5), Inches(9.2), Inches(1.0),
                 font_size=32, bold=True, color=TEXT_MAIN)

    # 4. Bullet points matching the '▸' style
    top = Inches(2.6)
    for point in points[:5]:
        # Small gold triangle/arrow
        _add_textbox(slide, "▸", 
                     Inches(0.4), top, Inches(0.3), Inches(0.5),
                     font_size=16, color=ACCENT)
        
        # Point text
        _add_textbox(slide, point,
                     Inches(0.7), top + Inches(0.02), Inches(8.5), Inches(0.8),
                     font_size=18, color=TEXT_MUTE)
        top += Inches(0.85)


# ──────────────────────────────────────────────────────────────────────────────
# MCP TOOLS — decorated with @mcp.tool()
# ──────────────────────────────────────────────────────────────────────────────

@mcp.tool()
def create_presentation(output_path: str) -> str:
    """
    Initialise a blank presentation.
    MUST be called before any other tool.

    Args:
        output_path: File path to save the .pptx e.g. 'workspace/my_deck.pptx'
    """
    global _prs, _output_path
    import os

    # ── Sanitise output_path ──────────────────────────────────────────────────
    output_path = (output_path or "").strip()

    # If empty or no .pptx extension → use safe default
    if not output_path:
        output_path = "presentation.pptx"
    if not output_path.endswith(".pptx"):
        output_path = output_path + ".pptx"

    # ── Always resolve to absolute path ──────────────────────────────────────
    _output_path = os.path.abspath(output_path)

    # ── Pre-create the parent folder ─────────────────────────────────────────
    parent = os.path.dirname(_output_path)
    if parent:
        os.makedirs(parent, exist_ok=True)

    _prs = Presentation()
    _prs.slide_width  = Inches(10)
    _prs.slide_height = Inches(7.5)
    return f"✅ Presentation initialised. Will save to '{_output_path}'."


@mcp.tool()
def add_slide(slide_type: str,
              title: str,
              subtitle: str = "",
              bullets: List = [],
              image_label: str = "[Image: Relevant Visual]") -> str:
    """
    Add a new slide to the presentation.

    Args:
        slide_type  : 'title'   → opening slide with big title + subtitle
                      'content' → title + 3-5 bullet points
                      'summary'  → closing slide with key-point cards
        title       : The slide title text.
        subtitle    : Optional subtitle (only for 'title' slides).
        bullets     : List of bullet strings (for 'content' / 'summary' slides).
        image_label : Label for the automatic image placeholder on content slides.
    """
    global _prs
    if _prs is None:
        return "❌ Call create_presentation first."

    blank_layout = _prs.slide_layouts[6]
    slide = _prs.slides.add_slide(blank_layout)

    idx = len(_prs.slides) - 1
    if slide_type == "title":
        _draw_title_slide(slide, title, subtitle, idx)
    elif slide_type == "summary":
        _draw_summary_slide(slide, title, bullets, idx)
    else:
        _draw_content_slide(slide, title, bullets, idx, image_label)

    return f"✅ Slide {idx} added ({slide_type}): '{title}'"


@mcp.tool()
def write_text(slide_index: int,
               title: str = "",
               bullets: List = [],
               image_label: str = "[Image: Relevant Visual]") -> str:
    """
    Overwrite the title and bullet text on an existing slide.

    Args:
        slide_index : 0-based index of the slide to update.
        title       : New title text.
        bullets     : New list of bullet point strings.
        image_label : New label for the image placeholder.
    """
    global _prs
    if _prs is None:
        return "❌ Call create_presentation first."
    if slide_index >= len(_prs.slides):
        return f"❌ Slide {slide_index} does not exist. Current count: {len(_prs.slides)}"

    slide = _prs.slides[slide_index]

    # Remove existing text boxes, keep background shapes
    sp_tree = slide.shapes._spTree
    to_remove = [sp._element for sp in slide.shapes if sp.has_text_frame]
    for el in to_remove:
        sp_tree.remove(el)

    if slide_index == 0:
        _draw_title_slide(slide, title, "", slide_index)
    elif slide_index == len(_prs.slides) - 1:
        _draw_summary_slide(slide, title, bullets, slide_index)
    else:
        _draw_content_slide(slide, title, bullets, slide_index, image_label)
    return f"✅ Slide {slide_index} text updated."


@mcp.tool()
def add_image_placeholder(slide_index: int,
                          label: str,
                          position: str = "right") -> str:
    """
    Add a labelled rectangle placeholder on a slide to represent an image.

    Args:
        slide_index : 0-based index of the target slide.
        label       : Descriptive label e.g. '[Image: Stellar nebula]'
        position    : 'left', 'right', or 'center'.
    """
    global _prs
    if _prs is None:
        return "❌ Call create_presentation first."
    if slide_index >= len(_prs.slides):
        return f"❌ Slide {slide_index} does not exist."

    slide = _prs.slides[slide_index]

    if position == "left":
        left = Inches(0.4)
    elif position == "center":
        left = Inches(2.5)
    else:
        left = Inches(5.5)

    ph = slide.shapes.add_shape(1, left, Inches(1.5), Inches(4.0), Inches(4.5))
    ph.fill.solid()
    ph.fill.fore_color.rgb = RGBColor(0x20, 0x30, 0x45)
    ph.line.color.rgb = ACCENT
    ph.line.width = Pt(1.5)

    _add_textbox(slide, label,
                 left + Inches(0.1), Inches(3.3),
                 Inches(3.8), Inches(0.7),
                 font_size=13, color=TEXT_MUTE,
                 italic=True, align=PP_ALIGN.CENTER)

    return f"✅ Image placeholder '{label}' added to slide {slide_index} ({position})."


@mcp.tool()
def save_presentation() -> str:
    """
    Save the presentation to disk.
    MUST be called as the final step after all slides are added.
    """
    import os
    global _prs, _output_path

    if _prs is None:
        return "❌ No presentation to save. Call create_presentation first."

    # ── Sanitise path ─────────────────────────────────────────────────────────
    path = (_output_path or "").strip()
    if not path:
        path = "presentation.pptx"
    if not path.endswith(".pptx"):
        path = path + ".pptx"

    # ── Always use absolute path ──────────────────────────────────────────────
    abs_path = os.path.abspath(path)

    # ── Safely create parent folder ───────────────────────────────────────────
    parent = os.path.dirname(abs_path)
    if parent and not os.path.exists(parent):
        os.makedirs(parent, exist_ok=True)

    # ── Save ──────────────────────────────────────────────────────────────────
    _prs.save(abs_path)
    return f"✅ Presentation saved to '{abs_path}' with {len(_prs.slides)} slides."


@mcp.tool()
def get_slide_count() -> str:
    """
    Return the number of slides currently in the presentation.
    Useful for the agent to verify its progress.
    """
    if _prs is None:
        return "0"
    return str(len(_prs.slides))


# ──────────────────────────────────────────────────────────────────────────────
# Entry point
# ──────────────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    mcp.run()