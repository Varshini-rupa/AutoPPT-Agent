"""
app.py  –  Streamlit Frontend for Auto-PPT Agent
"""

import json
import logging
import os
import sys
import time
import traceback
from datetime import datetime

import streamlit as st
from dotenv import load_dotenv

load_dotenv()

# ── Page config ────────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="Auto-PPT Agent",
    page_icon="✦",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ── CSS ────────────────────────────────────────────────────────────────────────
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=DM+Serif+Display:ital@0;1&family=DM+Sans:wght@300;400;500;600&display=swap');

html, body, [class*="css"] {
    font-family: 'DM Sans', sans-serif;
    background-color: #F7F6F2;
    color: #1A1A1A;
}
#MainMenu, footer, header { visibility: hidden; }
.block-container { padding-top: 2rem; padding-bottom: 2rem; }

.app-header { display:flex; align-items:center; gap:14px; margin-bottom:2rem; }
.app-logo {
    width:44px; height:44px; background:#1A1A1A; border-radius:12px;
    display:flex; align-items:center; justify-content:center;
    font-size:20px; flex-shrink:0; color:white;
}
.app-title {
    font-family:'DM Serif Display',serif; font-size:2rem;
    font-weight:400; color:#1A1A1A; line-height:1; margin:0;
}
.app-subtitle { font-size:0.82rem; color:#888; font-weight:300; margin:0; }

.stTextArea textarea {
    font-family:'DM Sans',sans-serif !important;
    font-size:0.95rem !important;
    border:1.5px solid #E0DDD6 !important;
    border-radius:12px !important;
    background:#FFFFFF !important;
    color:#1A1A1A !important;
    padding:14px 16px !important;
    resize:none !important;
}
.stTextArea textarea:focus { border-color:#1A1A1A !important; box-shadow:none !important; }
.stTextArea textarea::placeholder { color:#BBBBB0 !important; }

.stButton > button {
    font-family:'DM Sans',sans-serif !important; font-weight:500 !important;
    font-size:0.9rem !important; border-radius:10px !important;
    border:none !important; padding:0.6rem 1.4rem !important;
    transition:all 0.18s ease !important; cursor:pointer !important;
}
.stButton > button[kind="primary"] { background:#1A1A1A !important; color:#FFFFFF !important; }
.stButton > button[kind="primary"]:hover { background:#333 !important; transform:translateY(-1px); }
.stButton > button[kind="secondary"] {
    background:#F0EDE8 !important; color:#1A1A1A !important;
    border:1.5px solid #E0DDD6 !important;
}
.stButton > button[kind="secondary"]:hover { background:#E8E4DE !important; }

.stDownloadButton > button {
    font-family:'DM Sans',sans-serif !important; font-weight:600 !important;
    font-size:0.95rem !important; background:#FFD700 !important;
    color:#1A1A1A !important; border:none !important;
    border-radius:10px !important; padding:0.65rem 1.8rem !important;
    width:100% !important; transition:all 0.18s ease !important;
}
.stDownloadButton > button:hover {
    background:#F5CC00 !important; transform:translateY(-1px);
    box-shadow:0 4px 14px rgba(255,215,0,0.4) !important;
}

.log-container {
    background:#1A1A1A; border-radius:12px; padding:1rem 1.2rem;
    font-family:'Courier New',monospace; font-size:0.78rem;
    max-height:340px; overflow-y:auto; line-height:1.7;
}
.log-thought { color:#7ECFFF; }
.log-action  { color:#FFD700; }
.log-obs     { color:#98FF98; }
.log-error   { color:#FF8080; }
.log-info    { color:#AAAAAA; }
.log-finish  { color:#FFD700; font-weight:bold; }

.slide-card {
    background:#0D1B2A; border-radius:12px; padding:1.2rem 1.4rem;
    margin-bottom:0.7rem; border:1px solid #1B263B;
    position:relative; overflow:hidden;
}
.slide-card::before {
    content:''; position:absolute; left:0; top:0; bottom:0;
    width:3px; background:#FFD700; border-radius:3px 0 0 3px;
}
.slide-num  { font-size:0.7rem; color:#CADCFC; font-weight:500; letter-spacing:0.08em; text-transform:uppercase; margin-bottom:0.3rem; }
.slide-title { font-family:'DM Serif Display',serif; font-size:1.05rem; color:#FFFFFF; margin-bottom:0.5rem; }
.slide-bullet { font-size:0.8rem; color:#CADCFC; margin:0.2rem 0; padding-left:0.8rem; position:relative; }
.slide-bullet::before { content:'▸'; position:absolute; left:0; color:#FFD700; font-size:0.65rem; }
.slide-type-badge { display:inline-block; font-size:0.65rem; font-weight:600; letter-spacing:0.06em; text-transform:uppercase; padding:2px 8px; border-radius:20px; margin-bottom:0.5rem; }
.badge-title   { background:#FFD70022; color:#FFD700; }
.badge-content { background:#7ECFFF22; color:#7ECFFF; }
.badge-summary { background:#98FF9822; color:#98FF98; }

.section-label { font-size:0.72rem; font-weight:600; letter-spacing:0.1em; text-transform:uppercase; color:#AAAAAA; margin-bottom:0.6rem; margin-top:1.2rem; }

.sidebar-header { font-family:'DM Serif Display',serif; font-size:1.1rem; color:#1A1A1A; margin-bottom:1rem; padding-bottom:0.6rem; border-bottom:1.5px solid #E8E4DE; }
.history-item { background:#F7F6F2; border:1px solid #E8E4DE; border-radius:10px; padding:0.75rem 1rem; margin-bottom:0.5rem; }
.history-prompt { font-size:0.82rem; font-weight:500; color:#1A1A1A; white-space:nowrap; overflow:hidden; text-overflow:ellipsis; }
.history-time   { font-size:0.72rem; color:#AAAAAA; margin-top:2px; }
.history-status { font-size:0.7rem; margin-top:4px; }

/* popup */
.popup-overlay {
    position:fixed; inset:0; background:rgba(0,0,0,0.4); z-index:9999;
    display:flex; align-items:center; justify-content:center;
    animation:fadeIn 0.25s ease;
}
@keyframes fadeIn { from{opacity:0} to{opacity:1} }
.popup-box {
    background:#FFFFFF; border-radius:20px; padding:2.4rem 2.8rem;
    text-align:center; max-width:380px;
    box-shadow:0 20px 60px rgba(0,0,0,0.18);
    animation:slideUp 0.3s ease;
}
@keyframes slideUp { from{transform:translateY(30px);opacity:0} to{transform:translateY(0);opacity:1} }
.popup-icon  { font-size:3rem; margin-bottom:0.6rem; }
.popup-title { font-family:'DM Serif Display',serif; font-size:1.5rem; color:#1A1A1A; margin-bottom:0.4rem; }
.popup-msg   { font-size:0.88rem; color:#666; margin-bottom:0.5rem; }
.popup-path  { font-size:0.78rem; color:#888; font-family:'Courier New',monospace; background:#F7F6F2; padding:6px 12px; border-radius:6px; word-break:break-all; }

.stProgress > div > div { background:#FFD700 !important; }
</style>
""", unsafe_allow_html=True)


# ── Session state ──────────────────────────────────────────────────────────────
for k, v in {
    "history":        [],
    "logs":           [],
    "running":        False,
    "show_popup":     False,
    "last_pptx_path": None,
    "last_slides":    [],
    "prompt_value":   "",       # ← controls the text area value
}.items():
    if k not in st.session_state:
        st.session_state[k] = v


# ── LangChain callback ─────────────────────────────────────────────────────────
try:
    from langchain_core.callbacks.base import BaseCallbackHandler

    class StreamlitLogCallback(BaseCallbackHandler):
        def _log(self, tag, text):
            st.session_state["logs"].append({"tag": tag, "text": text.strip()})

        def on_llm_start(self, serialized, prompts, **kwargs):
            self._log("info", "🧠 LLM thinking...")

        def on_agent_action(self, action, **kwargs):
            self._log("action", f"⚡ Action: {action.tool}")
            self._log("action", f"   Input:  {str(action.tool_input)[:200]}")

        def on_tool_start(self, serialized, input_str, **kwargs):
            self._log("action", f"🔧 Running tool: {serialized.get('name','tool')}")

        def on_tool_end(self, output, **kwargs):
            preview = str(output)[:180]
            self._log("obs", f"👁  {preview}{'...' if len(str(output))>180 else ''}")

        def on_tool_error(self, error, **kwargs):
            self._log("error", f"❌ Tool error: {error}")

        def on_agent_finish(self, finish, **kwargs):
            self._log("finish", f"✅ {str(finish.return_values.get('output',''))[:120]}")

        def on_llm_error(self, error, **kwargs):
            self._log("error", f"❌ LLM error: {error}")

except ImportError:
    StreamlitLogCallback = None
    st.error("LangChain core not installed correctly for logging.")


# ── Slide extractor ────────────────────────────────────────────────────────────
def extract_slides(pptx_path):
    try:
        from pptx import Presentation as PRS
        prs    = PRS(pptx_path)
        slides = []
        total  = len(prs.slides)
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        # Skip decorative elements like 'SLIDE X' or badges
                        if t and not t.startswith("SLIDE ") and t not in ("TITLE", "CONTENT", "SUMMARY", "▸"):
                            texts.append(t)
            
            title   = texts[0] if texts else f"Slide {i+1}"
            bullets = texts[1:] if len(texts) > 1 else []
            # Clean up the bullet markers if they were caught
            bullets = [b.lstrip("▸ ").strip() for b in bullets]
            
            stype   = "title" if i == 0 else ("summary" if i == total-1 else "content")
            slides.append({"index": i, "type": stype, "title": title, "bullets": bullets})
        return slides
    except Exception:
        return []


# ── Log renderer ───────────────────────────────────────────────────────────────
def render_logs(logs):
    tag_map = {"thought":"log-thought","action":"log-action","obs":"log-obs",
               "error":"log-error","info":"log-info","finish":"log-finish"}
    lines = []
    for e in logs:
        cls  = tag_map.get(e["tag"], "log-info")
        text = e["text"].replace("<","&lt;").replace(">","&gt;")
        lines.append(f'<span class="{cls}">{text}</span>')
    return '<div class="log-container">' + "<br>".join(lines) + "</div>"


# ── Agent runner ───────────────────────────────────────────────────────────────
# ── Agent runner ───────────────────────────────────────────────────────────────
def run_agent(user_request):
    try:
        # Re-load env just in case
        load_dotenv(override=True)
        google_key  = os.getenv("GOOGLE_API_KEY", "")
        model_id    = os.getenv("MODEL_ID",        "gemini-2.0-flash")
        max_tokens  = int(os.getenv("MAX_TOKENS",  "1024"))
        temperature = float(os.getenv("TEMPERATURE","0.2"))
        max_iter    = int(os.getenv("AGENT_MAX_ITERATIONS","20"))
        workspace   = os.getenv("AGENT_WORKSPACE", "./workspace")
        os.makedirs(workspace, exist_ok=True)

        if not google_key:
            return False, "GOOGLE_API_KEY not set in .env file."

        # Create a unique filename for this run
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_name = "".join([c if c.isalnum() else "_" for c in user_request[:20]]).strip("_")
        output_file = os.path.join(workspace, f"ppt_{timestamp}_{safe_name}.pptx")

        from langchain.agents import AgentExecutor, create_react_agent
        from langchain.tools import Tool
        from langchain_core.prompts import PromptTemplate
        from langchain_google_genai import ChatGoogleGenerativeAI
        import pptx_mcp_server
        from search_mcp_server import search_topic

        def _p(s):
            """Safely parse agent's JSON string argument into a Python dict."""
            s = (str(s) or "").strip()
            if not s or s in ("{}", ""):
                return {}
            # If it's just a path string (legacy), wrap it
            if s.endswith(".pptx") and "{" not in s:
                return {"output_path": s}
            try:
                return json.loads(s)
            except Exception:
                try:
                    # Try fixing single quotes
                    return json.loads(s.replace("'",'"'))
                except Exception:
                    try:
                        # Fallback to eval (dangerous but sometimes needed for LLM hallucinations)
                        return eval(s)
                    except Exception:
                        return {}

        def t_create(a):
            x = _p(a)
            # Use OUR unique file path if they don't specify one or if it's generic
            p = x.get("output_path", output_file)
            if not p.endswith(".pptx"): p += ".pptx"
            if not os.path.isabs(p) and not p.startswith(workspace):
                p = os.path.join(workspace, os.path.basename(p))
            return pptx_mcp_server.create_presentation(p)

        def t_slide(a):
            x = _p(a)
            return pptx_mcp_server.add_slide(
                x.get("slide_type","content"), x.get("title","Slide"),
                x.get("subtitle",""), x.get("bullets",[])
            )

        def t_write(a):
            x = _p(a)
            return pptx_mcp_server.write_text(int(x.get("slide_index",0)), x.get("title",""), x.get("bullets",[]))

        def t_img(a):
            x = _p(a)
            return pptx_mcp_server.add_image_placeholder(
                int(x.get("slide_index",0)),
                x.get("label","[Image]"), x.get("position","right")
            )

        def t_save(a):
            return pptx_mcp_server.save_presentation()

        def t_count(a):
            return pptx_mcp_server.get_slide_count()

        def t_search(q):
            # Agent might pass JSON to search, but search_topic wants plain text
            if isinstance(q, str) and (q.startswith("{") or q.startswith("[")):
                try:
                    payload = json.loads(q)
                    if isinstance(payload, dict):
                        q = payload.get("query", next(iter(payload.values())))
                except: pass
            return search_topic(str(q).strip())

        tools = [
            Tool(name="create_presentation", func=t_create,
                 description='Initialise .pptx. Input JSON: {"output_path":"workspace/file.pptx"}'),
            Tool(name="add_slide", func=t_slide,
                 description='Add slide. Input JSON: {"slide_type":"title|content|summary","title":"...","subtitle":"...","bullets":["..."]}'),
            Tool(name="write_text", func=t_write,
                 description='Overwrite slide. Input JSON: {"slide_index":0,"title":"...","bullets":["..."]}'),
            Tool(name="add_image_placeholder", func=t_img,
                 description='Image placeholder. Input JSON: {"slide_index":1,"label":"[Image:...]","position":"right"}'),
            Tool(name="save_presentation", func=t_save,
                 description="Save file. Call LAST. Input: {}"),
            Tool(name="get_slide_count", func=t_count,
                 description="Slide count. Input: {}"),
            Tool(name="search_topic", func=t_search,
                 description="Search web for facts. Input: plain text search query."),
        ]

        # Use the more comprehensive template from agent_ppt.py
        PROMPT_TEMPLATE = """You are Auto-PPT, an expert presentation agent.
Your job: turn a user's single-sentence request into a polished PowerPoint file.

User Input -> Agent Brain (LLM) -> Decision: "Write Slide 1"
             -> Tool Call (MCP: write_slide)
             -> Agent Brain -> Decision: "Write Slide 2"
             -> Tool Call...
             -> Agent Brain -> Decision: "Save File" -> FINISH

MANDATORY 4-STEP AGENTIC LOOP:

STEP 1 - PLAN (MUST be done before ANY tool call)
  Think through: audience, tone, number of slides requested (default 5 if not specified).
  Produce a JSON outline in your Thought before tools.

STEP 2 - RESEARCH
  For each content slide topic, call search_topic once to gather real facts.

STEP 3 - BUILD
  a. Call create_presentation first.
  b. For each slide call add_slide with correct type, title, bullets.
     - title slides   : {{"slide_type":"title",   "title":"...", "subtitle":"..."}}
     - content slides : {{"slide_type":"content", "title":"...", "bullets":["...","...","..."]}}
     - summary slides : {{"slide_type":"summary", "title":"...", "bullets":["...","..."]}}

STEP 4 - SAVE
  Call save_presentation {{}}.
  Then reply: "Done! Presentation saved to {path}."

RULES:
- NEVER skip planning.
- ALWAYS have at least 1 title slide and 1 summary slide.
- Tools MUST receive valid JSON strings (except search_topic).
- If a tool errors, fix inputs and retry.

Tools: {tools}
Tool names: {tool_names}

Format:
Thought: <reasoning>
Action: <tool name>
Action Input: <JSON string>
Observation: <result>
... repeat ...
Thought: I now know the final answer.
Final Answer: <message>

Begin!
Question: {input}
{agent_scratchpad}"""

        llm = ChatGoogleGenerativeAI(
            model=model_id, temperature=temperature,
            max_output_tokens=max_tokens, google_api_key=google_key,
        )
        callbacks = [StreamlitLogCallback()] if StreamlitLogCallback else []
        prompt    = PromptTemplate.from_template(PROMPT_TEMPLATE.replace("{path}", output_file))
        agent     = create_react_agent(llm=llm, tools=tools, prompt=prompt)
        executor  = AgentExecutor(
            agent=agent, tools=tools, verbose=True, # Set to True for better server-side debugging
            max_iterations=max_iter, handle_parsing_errors=True,
            callbacks=callbacks,
        )

        st.session_state["logs"].append({"tag":"info","text":"🤖 Thinking about your presentation..."})
        executor.invoke({"input": user_request})

        # Check if the file was actually saved at our expected location
        saved_path = pptx_mcp_server._output_path
        if saved_path and os.path.exists(saved_path):
            return True, saved_path

        # Fallback to search if _output_path failed
        if os.path.exists(output_file):
            return True, output_file

        return False, "Agent finished but no .pptx found. Please check logs for errors."

    except Exception as e:
        return False, f"{type(e).__name__}: {e}\n{traceback.format_exc()}"


# ══════════════════════════════════════════════════════════════════════════════
# UI LAYOUT
# ══════════════════════════════════════════════════════════════════════════════

# ── Sidebar ────────────────────────────────────────────────────────────────────
with st.sidebar:
    st.markdown('<div class="sidebar-header">✦ History</div>', unsafe_allow_html=True)
    if not st.session_state["history"]:
        st.markdown('<p style="font-size:0.82rem;color:#AAAAAA;">No presentations yet.<br>Generate one to see it here.</p>',
                    unsafe_allow_html=True)
    else:
        for item in reversed(st.session_state["history"]):
            icon = "✅" if item["status"]=="done" else ("❌" if item["status"]=="error" else "⏳")
            st.markdown(f"""
            <div class="history-item">
                <div class="history-prompt">{item['prompt'][:55]}{'…' if len(item['prompt'])>55 else ''}</div>
                <div class="history-time">{item['time']}</div>
                <div class="history-status">{icon} {item['status'].capitalize()}</div>
            </div>""", unsafe_allow_html=True)
    st.markdown("---")
    st.markdown('<p style="font-size:0.75rem;color:#BBBBBB;">Auto-PPT Agent v1.0<br>Powered by Gemini + MCP</p>',
                unsafe_allow_html=True)


# ── Header ─────────────────────────────────────────────────────────────────────
st.markdown("""
<div class="app-header">
    <div class="app-logo">✦</div>
    <div>
        <div class="app-title">Auto-PPT Agent</div>
        <div class="app-subtitle">AI-powered presentations via MCP + Gemini</div>
    </div>
</div>
""", unsafe_allow_html=True)

# ── API key check ──────────────────────────────────────────────────────────────
if not os.getenv("GOOGLE_API_KEY"):
    st.error("⚠️ **GOOGLE_API_KEY** not found in your `.env` file. Get a free key at [aistudio.google.com](https://aistudio.google.com/apikey)")
    st.stop()

# ── Two columns ────────────────────────────────────────────────────────────────
col_left, col_right = st.columns([1.1, 0.9], gap="large")

with col_left:
    st.markdown('<div class="section-label">Your Prompt</div>', unsafe_allow_html=True)

    # ── Example chips ──────────────────────────────────────────────────────────
    examples = [
        "Create a 5-slide presentation on the life cycle of a star for a 6th-grade class",
        "Make a 6-slide deck on climate change for a university seminar",
        "Create a 4-slide presentation on the French Revolution for high school",
        "Make a 5-slide business pitch deck for a mobile app startup",
    ]

    chip_cols = st.columns(2)
    for idx, ex in enumerate(examples):
        with chip_cols[idx % 2]:
            if st.button(ex[:42]+"…", key=f"chip_{idx}", type="secondary",
                         disabled=st.session_state["running"]):
                # ✅ FIX: set prompt_value then rerun so textarea picks it up
                st.session_state["prompt_value"] = ex
                st.rerun()

    # ── Text area ──────────────────────────────────────────────────────────────
    # ✅ FIX: use value= instead of key= to avoid NoneType error
    user_prompt = st.text_area(
        label="prompt",
        label_visibility="collapsed",
        placeholder="e.g. Create a 5-slide presentation on the life cycle of a star…",
        height=100,
        value=st.session_state["prompt_value"],
        disabled=st.session_state["running"],
    )
    # Keep session state in sync
    st.session_state["prompt_value"] = user_prompt or ""

    # ── Buttons ────────────────────────────────────────────────────────────────
    b1, b2 = st.columns([2, 1])
    with b1:
        generate_clicked = st.button(
            "✦ Generate Presentation", type="primary",
            disabled=st.session_state["running"] or not (user_prompt or "").strip(),
            use_container_width=True,
        )
    with b2:
        if st.button("Clear", type="secondary", use_container_width=True,
                     disabled=st.session_state["running"]):
            st.session_state["logs"]           = []
            st.session_state["last_pptx_path"] = None
            st.session_state["last_slides"]    = []
            st.session_state["prompt_value"]   = ""
            st.rerun()

    # ── Log viewer ─────────────────────────────────────────────────────────────
    st.markdown('<div class="section-label" style="margin-top:1.4rem;">Agent Log</div>',
                unsafe_allow_html=True)
    log_ph = st.empty()
    if st.session_state["logs"]:
        log_ph.markdown(render_logs(st.session_state["logs"]), unsafe_allow_html=True)
    else:
        log_ph.markdown(
            '<div class="log-container"><span class="log-info">Waiting for agent to start…</span></div>',
            unsafe_allow_html=True)


with col_right:
    st.markdown('<div class="section-label">Slide Preview</div>', unsafe_allow_html=True)

    preview_ph  = st.empty()
    download_ph = st.empty()

    if st.session_state["last_slides"]:
        with preview_ph.container():
            for slide in st.session_state["last_slides"]:
                bc = f"badge-{slide['type']}"
                bullets_html = "".join(
                    f'<div class="slide-bullet">{b}</div>' for b in slide["bullets"][:4]
                ) or '<div class="slide-bullet" style="opacity:0.4">No bullets</div>'
                st.markdown(f"""
                <div class="slide-card">
                    <div class="slide-num">Slide {slide['index']+1}</div>
                    <span class="slide-type-badge {bc}">{slide['type']}</span>
                    <div class="slide-title">{slide['title']}</div>
                    {bullets_html}
                </div>""", unsafe_allow_html=True)
    else:
        preview_ph.markdown("""
        <div style="background:#F7F6F2;border:1.5px dashed #D4D0C8;border-radius:14px;
                    text-align:center;padding:3rem 1rem;">
            <div style="font-size:2.5rem;margin-bottom:0.6rem;">🎞️</div>
            <div style="font-size:0.88rem;color:#AAAAAA;">Slide preview will appear here<br>after generation.</div>
        </div>""", unsafe_allow_html=True)

    # ── Download button ────────────────────────────────────────────────────────
    if st.session_state["last_pptx_path"] and os.path.exists(st.session_state["last_pptx_path"]):
        with open(st.session_state["last_pptx_path"], "rb") as f:
            pptx_bytes = f.read()
        fname = os.path.basename(st.session_state["last_pptx_path"])
        with download_ph.container():
            st.markdown('<div style="margin-top:1rem;"></div>', unsafe_allow_html=True)
            if st.download_button(
                label="⬇ Download Presentation",
                data=pptx_bytes,
                file_name=fname,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
                key="dl_btn",
            ):
                st.session_state["show_popup"] = True
                st.rerun()


# ── Success popup ──────────────────────────────────────────────────────────────
if st.session_state["show_popup"]:
    pptx_abs = os.path.abspath(st.session_state["last_pptx_path"] or "")
    fname    = os.path.basename(pptx_abs)
    st.markdown(f"""
    <div class="popup-overlay">
        <div class="popup-box">
            <div class="popup-icon">🎉</div>
            <div class="popup-title">Download Started!</div>
            <div class="popup-msg">Your presentation is downloading to your browser's Downloads folder.</div>
            <div class="popup-path">Also saved in workspace:<br>{pptx_abs}</div>
        </div>
    </div>
    """, unsafe_allow_html=True)
    time.sleep(3)
    st.session_state["show_popup"] = False
    st.rerun()


# ══════════════════════════════════════════════════════════════════════════════
# RUN AGENT
# ══════════════════════════════════════════════════════════════════════════════
if generate_clicked and (user_prompt or "").strip():
    st.session_state["running"]        = True
    st.session_state["logs"]           = []
    st.session_state["last_pptx_path"] = None
    st.session_state["last_slides"]    = []

    st.session_state["history"].append({
        "prompt":    (user_prompt or "").strip(),
        "time":      datetime.now().strftime("%d %b %Y, %H:%M"),
        "status":    "running",
        "pptx_path": None,
    })

    st.session_state["logs"].append({"tag":"info","text":f"▶ Starting: {user_prompt[:80]}…"})
    log_ph.markdown(render_logs(st.session_state["logs"]), unsafe_allow_html=True)

    with st.spinner("Agent is working…"):
        success, result = run_agent((user_prompt or "").strip())

    st.session_state["running"] = False

    if success:
        slides = extract_slides(result)
        st.session_state["last_pptx_path"] = result
        st.session_state["last_slides"]    = slides
        st.session_state["logs"].append({"tag":"finish","text":f"✅ Saved: {os.path.basename(result)}"})
        st.session_state["history"][-1]["status"] = "done"
        # ✅ Also show a Streamlit success banner
        st.success(f"✅ Presentation created! **{os.path.basename(result)}** is ready to download.")
        st.balloons()
    else:
        st.session_state["logs"].append({"tag":"error","text":f"❌ {result[:300]}"})
        st.session_state["history"][-1]["status"] = "error"
        st.error(f"❌ Agent failed: {result[:300]}")

    st.rerun()